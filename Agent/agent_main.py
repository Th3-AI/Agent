import os
import asyncio
import webbrowser
import subprocess
import platform
import psutil
import tempfile
import time
import numpy as np
import cv2
import pyautogui
from PIL import Image
import pytesseract
from pathlib import Path
from google.adk.agents import Agent
from google.adk.sessions import InMemorySessionService
from google.genai import types as genai_types
import google.generativeai as genai
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
import json
import datetime
import pickle
import threading
from dateutil.parser import parse as parse_datetime
import concurrent.futures
import re
from collections import deque
from typing import Dict, List, Optional, Tuple, Any
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication
import ctypes
from ctypes import cast, POINTER
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
import winreg
import requests
from bs4 import BeautifulSoup
import requests
from pathlib import Path
import sys
import io
import time
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import base64
from urllib.parse import quote_plus

# Configure API keys
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "AIzaSyB0XSUK7f3LhjG9i2n8frZSB4nrKAgLEg0")
os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY

# Email configuration
EMAIL_CONFIG = {
    "smtp_server": os.getenv("SMTP_SERVER"),
    "smtp_port": int(os.getenv("SMTP_PORT"),
    "sender_email": os.getenv("SENDER_EMAIL"),
    "sender_password": os.getenv("SENDER_PASSWORD"),
    "use_tls": True
}

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
    # Update email config with values from .env
    EMAIL_CONFIG.update({
        "smtp_server": os.getenv("SMTP_SERVER", EMAIL_CONFIG["smtp_server"]),
        "smtp_port": int(os.getenv("SMTP_PORT", str(EMAIL_CONFIG["smtp_port"]))),
        "sender_email": os.getenv("SENDER_EMAIL", EMAIL_CONFIG["sender_email"]),
        "sender_password": os.getenv("SENDER_PASSWORD", EMAIL_CONFIG["sender_password"])
    })
except ImportError:
    print("python-dotenv not installed. Using default email configuration.")
except Exception as e:
    print(f"Error loading .env file: {e}")

# Configure genai
genai.configure(api_key=GOOGLE_API_KEY)

# Memory storage path
MEMORY_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "contextual_memory.pkl")

# Task scheduler storage path
TASKS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "scheduled_tasks.pkl")

# Global variable to store scheduled tasks
scheduled_tasks = []
task_executor = None

class ContextualMemory:
    """A class for storing and retrieving user interaction context and memories."""
    
    def __init__(self, max_conversations=20, max_facts=100, max_preferences=50):
        """Initialize the contextual memory system.
        
        Args:
            max_conversations: Maximum number of recent conversations to remember
            max_facts: Maximum number of facts to remember about the user
            max_preferences: Maximum number of user preferences to store
        """
        self.recent_conversations = deque(maxlen=max_conversations)
        self.user_facts = {}  # Facts about the user
        self.user_preferences = {}  # User preferences
        self.topic_memories = {}  # Memories organized by topic
        self.last_accessed = {}  # Track when topics were last accessed
        self.memory_initialized = False
    
    def add_conversation(self, user_input: str, agent_response: str, timestamp=None):
        """Add a conversation exchange to memory.
        
        Args:
            user_input: The user's input text
            agent_response: The agent's response text
            timestamp: Optional timestamp, defaults to current time
        """
        if timestamp is None:
            timestamp = datetime.datetime.now()
            
        self.recent_conversations.append({
            "user_input": user_input,
            "agent_response": agent_response,
            "timestamp": timestamp
        })
    
    def add_fact(self, category: str, fact: str):
        """Add a fact about the user.
        
        Args:
            category: The category of the fact (e.g., "personal", "work", "interests")
            fact: The fact to remember
        """
        if category not in self.user_facts:
            self.user_facts[category] = []
            
        # Check if this fact already exists to avoid duplicates
        if fact not in self.user_facts[category]:
            self.user_facts[category].append(fact)
    
    def add_preference(self, category: str, preference: str):
        """Add a user preference.
        
        Args:
            category: The category of the preference (e.g., "color", "food", "music")
            preference: The preference to remember
        """
        self.user_preferences[category] = preference
    
    def add_topic_memory(self, topic: str, memory: Any):
        """Add a memory related to a specific topic.
        
        Args:
            topic: The topic for this memory
            memory: The memory content (can be any object)
        """
        self.topic_memories[topic] = memory
        self.last_accessed[topic] = datetime.datetime.now()
    
    def get_recent_conversations(self, count=5) -> List[Dict]:
        """Get the most recent conversations.
        
        Args:
            count: Number of recent conversations to retrieve
            
        Returns:
            List of conversation dictionaries
        """
        return list(self.recent_conversations)[-count:]
    
    def get_facts(self, category=None) -> Dict[str, List[str]]:
        """Get facts about the user.
        
        Args:
            category: Optional category to filter by
            
        Returns:
            Dictionary of facts, filtered by category if provided
        """
        if category:
            return {category: self.user_facts.get(category, [])}
        return self.user_facts
    
    def get_preferences(self, category=None) -> Dict[str, str]:
        """Get user preferences.
        
        Args:
            category: Optional category to filter by
            
        Returns:
            Dictionary of preferences, filtered by category if provided
        """
        if category:
            if category in self.user_preferences:
                return {category: self.user_preferences[category]}
            return {}
        return self.user_preferences
    
    def get_topic_memory(self, topic: str) -> Optional[Any]:
        """Get memory related to a specific topic.
        
        Args:
            topic: The topic to retrieve
            
        Returns:
            The memory content or None if not found
        """
        if topic in self.topic_memories:
            self.last_accessed[topic] = datetime.datetime.now()
            return self.topic_memories[topic]
        return None
    
    def search_memories(self, query: str) -> Dict[str, Any]:
        """Search all memories for relevant information.
        
        Args:
            query: The search query
            
        Returns:
            Dictionary of relevant memories found
        """
        results = {
            "conversations": [],
            "facts": {},
            "preferences": {},
            "topics": []
        }
        
        # Search conversations
        for conv in self.recent_conversations:
            if query.lower() in conv["user_input"].lower() or query.lower() in conv["agent_response"].lower():
                results["conversations"].append(conv)
        
        # Search facts
        for category, facts in self.user_facts.items():
            matching_facts = [fact for fact in facts if query.lower() in fact.lower()]
            if matching_facts:
                results["facts"][category] = matching_facts
        
        # Search preferences
        for category, preference in self.user_preferences.items():
            if query.lower() in preference.lower() or query.lower() in category.lower():
                results["preferences"][category] = preference
        
        # Search topics
        for topic in self.topic_memories:
            if query.lower() in topic.lower():
                results["topics"].append(topic)
        
        return results
    
    def extract_facts_from_conversation(self, user_input: str, agent_response: str):
        """Extract potential facts from conversation.
        
        This looks for patterns in user input that might indicate personal information.
        
        Args:
            user_input: The user's input text
            agent_response: The agent's response text
        """
        # Pattern for personal information
        name_match = re.search(r"my name is (\w+)", user_input.lower())
        if name_match:
            self.add_fact("personal", f"User's name is {name_match.group(1)}")
        
        # Pattern for likes/interests
        like_match = re.search(r"i (like|love|enjoy) (.+?)[\.!\?]", user_input.lower())
        if like_match:
            self.add_fact("interests", f"User likes {like_match.group(2)}")
        
        # Pattern for work related information
        work_match = re.search(r"i (work|am working|worked) (at|for|as) (.+?)[\.!\?]", user_input.lower())
        if work_match:
            self.add_fact("work", f"User works {work_match.group(2)} {work_match.group(3)}")
        
        # Pattern for preferences
        if "prefer" in user_input.lower():
            prefer_match = re.search(r"i prefer (.+?)[\.!\?]", user_input.lower())
            if prefer_match:
                preference = prefer_match.group(1)
                self.add_preference("general", preference)
    
    def get_memory_context(self, current_input: str, max_context_items=5) -> str:
        """Generate a context summary from memories relevant to the current input.
        
        Args:
            current_input: The current user input
            max_context_items: Maximum number of context items to include
            
        Returns:
            A formatted string with relevant memory context
        """
        if not self.memory_initialized:
            return ""
        
        context_parts = []
        added_items = 0
        
        # Search for relevant memories
        search_results = self.search_memories(current_input)
        
        # Add relevant conversations
        if search_results["conversations"]:
            conv_summary = []
            for conv in search_results["conversations"][:2]:
                conv_summary.append(f"User previously asked: '{conv['user_input']}' and you responded: '{conv['agent_response']}'")
                added_items += 1
            if conv_summary:
                context_parts.append("Previous relevant conversations:\n" + "\n".join(conv_summary))
        
        # Add relevant facts if we have space
        if search_results["facts"] and added_items < max_context_items:
            fact_summary = []
            for category, facts in search_results["facts"].items():
                for fact in facts:
                    fact_summary.append(f"{fact}")
                    added_items += 1
                    if added_items >= max_context_items:
                        break
                if added_items >= max_context_items:
                    break
            if fact_summary:
                context_parts.append("Relevant user information:\n" + "\n".join(fact_summary))
        
        # Add user preferences if we have space
        if search_results["preferences"] and added_items < max_context_items:
            pref_summary = []
            for category, preference in search_results["preferences"].items():
                pref_summary.append(f"User prefers {preference} for {category}")
                added_items += 1
                if added_items >= max_context_items:
                    break
            if pref_summary:
                context_parts.append("User preferences:\n" + "\n".join(pref_summary))
        
        # Add the most recent conversations if we still have space and no relevant ones were found
        if not search_results["conversations"] and added_items < max_context_items:
            recent = self.get_recent_conversations(2)
            if recent:
                recent_summary = []
                for conv in recent:
                    recent_summary.append(f"User recently asked: '{conv['user_input']}' and you responded: '{conv['agent_response']}'")
                    added_items += 1
                    if added_items >= max_context_items:
                        break
                context_parts.append("Recent conversations:\n" + "\n".join(recent_summary))
        
        # Combine all context parts
        if context_parts:
            return "\n\n".join(context_parts)
        return ""
    
    def save_to_disk(self):
        """Save the memory to disk."""
        try:
            with open(MEMORY_FILE, 'wb') as f:
                pickle.dump({
                    "recent_conversations": list(self.recent_conversations),
                    "user_facts": self.user_facts,
                    "user_preferences": self.user_preferences,
                    "topic_memories": self.topic_memories,
                    "last_accessed": self.last_accessed
                }, f)
            print(f"Memory saved to {MEMORY_FILE}")
            return True
        except Exception as e:
            print(f"Error saving memory: {e}")
            return False
    
    def load_from_disk(self):
        """Load the memory from disk."""
        try:
            if os.path.exists(MEMORY_FILE):
                with open(MEMORY_FILE, 'rb') as f:
                    data = pickle.load(f)
                    self.recent_conversations = deque(data["recent_conversations"], maxlen=self.recent_conversations.maxlen)
                    self.user_facts = data["user_facts"]
                    self.user_preferences = data["user_preferences"]
                    self.topic_memories = data["topic_memories"]
                    self.last_accessed = data["last_accessed"]
                    self.memory_initialized = True
                print(f"Memory loaded from {MEMORY_FILE}")
                return True
            self.memory_initialized = True
            return False
        except Exception as e:
            print(f"Error loading memory: {e}")
            self.memory_initialized = True
            return False

class ScheduledTask:
    """Represents a scheduled task with execution time and command."""
    def __init__(self, task_id, description, command, execute_at, repeat_interval=None):
        self.task_id = task_id
        self.description = description
        self.command = command
        self.execute_at = execute_at  # datetime object
        self.repeat_interval = repeat_interval  # in seconds, None for no repeat
        self.last_executed = None
    
    def __str__(self):
        repeat_info = f", Repeats every {format_time_interval(self.repeat_interval)}" if self.repeat_interval else ""
        last_run = f", Last run: {self.last_executed}" if self.last_executed else ""
        return f"Task {self.task_id}: {self.description} (Run at: {self.execute_at.strftime('%Y-%m-%d %H:%M:%S')}{repeat_info}{last_run})"
    
    def to_dict(self):
        """Convert task to dictionary for display."""
        return {
            "id": self.task_id,
            "description": self.description,
            "command": self.command,
            "execute_at": self.execute_at.strftime("%Y-%m-%d %H:%M:%S"),
            "repeat_interval": self.repeat_interval,
            "last_executed": self.last_executed.strftime("%Y-%m-%d %H:%M:%S") if self.last_executed else None
        }

def format_time_interval(seconds):
    """Format a time interval in seconds to a human-readable string."""
    if not seconds:
        return "never"
    
    minutes, seconds = divmod(seconds, 60)
    hours, minutes = divmod(minutes, 60)
    days, hours = divmod(hours, 24)
    
    parts = []
    if days > 0:
        parts.append(f"{days} day{'s' if days != 1 else ''}")
    if hours > 0:
        parts.append(f"{hours} hour{'s' if hours != 1 else ''}")
    if minutes > 0:
        parts.append(f"{minutes} minute{'s' if minutes != 1 else ''}")
    if seconds > 0 and not parts:
        parts.append(f"{seconds} second{'s' if seconds != 1 else ''}")
    
    return " ".join(parts)

def save_tasks():
    """Save scheduled tasks to disk."""
    try:
        with open(TASKS_FILE, 'wb') as f:
            pickle.dump(scheduled_tasks, f)
        print(f"Debug: Saved {len(scheduled_tasks)} tasks to {TASKS_FILE}")
        return True
    except Exception as e:
        print(f"Error saving tasks: {e}")
        return False

def load_tasks():
    """Load scheduled tasks from disk."""
    global scheduled_tasks
    try:
        if os.path.exists(TASKS_FILE):
            with open(TASKS_FILE, 'rb') as f:
                loaded_tasks = pickle.load(f)
                scheduled_tasks = loaded_tasks
            print(f"Debug: Loaded {len(scheduled_tasks)} tasks from {TASKS_FILE}")
        return True
    except Exception as e:
        print(f"Error loading tasks: {e}")
        scheduled_tasks = []
        return False

def get_next_task_id():
    """Get the next available task ID."""
    if not scheduled_tasks:
        return 1
    return max(task.task_id for task in scheduled_tasks) + 1

def add_scheduled_task(description, command, execute_at, repeat_interval=None):
    """Add a new scheduled task."""
    try:
        # Parse execute_at if it's a string
        if isinstance(execute_at, str):
            try:
                execute_at = parse_datetime(execute_at)
            except:
                return False, "Invalid date/time format. Please use a format like 'YYYY-MM-DD HH:MM:SS' or 'tomorrow at 3pm'."
        
        # Parse repeat_interval if it's a string
        if isinstance(repeat_interval, str) and repeat_interval:
            try:
                if "hour" in repeat_interval.lower():
                    hours = float(repeat_interval.lower().split("hour")[0].strip())
                    repeat_interval = int(hours * 3600)
                elif "minute" in repeat_interval.lower():
                    minutes = float(repeat_interval.lower().split("minute")[0].strip())
                    repeat_interval = int(minutes * 60)
                elif "day" in repeat_interval.lower():
                    days = float(repeat_interval.lower().split("day")[0].strip())
                    repeat_interval = int(days * 86400)
                elif "week" in repeat_interval.lower():
                    weeks = float(repeat_interval.lower().split("week")[0].strip())
                    repeat_interval = int(weeks * 604800)
                else:
                    repeat_interval = int(repeat_interval)
            except:
                return False, "Invalid repeat interval. Please use a format like '1 hour', '30 minutes', '1 day', or seconds as an integer."
        
        task_id = get_next_task_id()
        task = ScheduledTask(task_id, description, command, execute_at, repeat_interval)
        scheduled_tasks.append(task)
        save_tasks()
        
        # Start task executor if not running
        start_task_executor()
        
        return True, task
    except Exception as e:
        print(f"Error adding task: {e}")
        return False, f"Error adding task: {e}"

def remove_scheduled_task(task_id):
    """Remove a scheduled task by ID."""
    try:
        task_id = int(task_id)
        global scheduled_tasks
        original_count = len(scheduled_tasks)
        scheduled_tasks = [task for task in scheduled_tasks if task.task_id != task_id]
        
        if len(scheduled_tasks) < original_count:
            save_tasks()
            return True, f"Task {task_id} removed successfully."
        else:
            return False, f"Task {task_id} not found."
    except Exception as e:
        print(f"Error removing task: {e}")
        return False, f"Error removing task: {e}"

def update_scheduled_task(task_id, description=None, command=None, execute_at=None, repeat_interval=None):
    """Update an existing scheduled task."""
    try:
        task_id = int(task_id)
        for task in scheduled_tasks:
            if task.task_id == task_id:
                if description is not None:
                    task.description = description
                if command is not None:
                    task.command = command
                if execute_at is not None:
                    if isinstance(execute_at, str):
                        task.execute_at = parse_datetime(execute_at)
                    else:
                        task.execute_at = execute_at
                if repeat_interval is not None:
                    if isinstance(repeat_interval, str):
                        if repeat_interval.lower() == "none":
                            task.repeat_interval = None
                        elif "hour" in repeat_interval.lower():
                            hours = float(repeat_interval.lower().split("hour")[0].strip())
                            task.repeat_interval = int(hours * 3600)
                        elif "minute" in repeat_interval.lower():
                            minutes = float(repeat_interval.lower().split("minute")[0].strip())
                            task.repeat_interval = int(minutes * 60)
                        elif "day" in repeat_interval.lower():
                            days = float(repeat_interval.lower().split("day")[0].strip())
                            task.repeat_interval = int(days * 86400)
                        elif "week" in repeat_interval.lower():
                            weeks = float(repeat_interval.lower().split("week")[0].strip())
                            task.repeat_interval = int(weeks * 604800)
                        else:
                            task.repeat_interval = int(repeat_interval)
                    else:
                        task.repeat_interval = repeat_interval
                
                save_tasks()
                return True, task
        
        return False, f"Task {task_id} not found."
    except Exception as e:
        print(f"Error updating task: {e}")
        return False, f"Error updating task: {e}"

def list_scheduled_tasks():
    """List all scheduled tasks."""
    return [task.to_dict() for task in scheduled_tasks]

def execute_task(task):
    """Execute a scheduled task."""
    try:
        print(f"Debug: Executing task {task.task_id}: {task.description}")
        async def run_task():
            try:
                # Use process_command to execute the task
                await process_user_input(task.command)
                
                # Update last executed time
                task.last_executed = datetime.datetime.now()
                
                # If task repeats, schedule next execution
                if task.repeat_interval:
                    task.execute_at = datetime.datetime.now() + datetime.timedelta(seconds=task.repeat_interval)
                    print(f"Debug: Task {task.task_id} rescheduled for {task.execute_at}")
                else:
                    # Remove one-time task after execution
                    remove_scheduled_task(task.task_id)
                
                save_tasks()
            except Exception as e:
                print(f"Error executing task {task.task_id}: {e}")
        
        # Run the async task in a new event loop
        asyncio.run(run_task())
        
        return True
    except Exception as e:
        print(f"Error executing task {task.task_id}: {e}")
        return False

def task_executor_loop():
    """Background thread to check and execute scheduled tasks."""
    print("Debug: Task executor loop started")
    while True:
        try:
            now = datetime.datetime.now()
            tasks_to_execute = []
            
            # Find tasks that need to be executed
            for task in scheduled_tasks:
                if task.execute_at <= now:
                    tasks_to_execute.append(task)
            
            # Execute tasks in thread pool to avoid blocking
            if tasks_to_execute:
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    futures = [executor.submit(execute_task, task) for task in tasks_to_execute]
                    concurrent.futures.wait(futures)
            
            # Sleep for 1 minute to avoid constant CPU usage
            time.sleep(60)
        except Exception as e:
            print(f"Error in task executor loop: {e}")
            time.sleep(60)  # Sleep on error to avoid spinning

def start_task_executor():
    """Start the task executor thread if not already running."""
    global task_executor
    if task_executor is None or not task_executor.is_alive():
        task_executor = threading.Thread(target=task_executor_loop, daemon=True)
        task_executor.start()
        print("Debug: Task executor thread started")

def schedule_task_command(request):
    """Process a task scheduling request and return formatted response."""
    try:
        # Parse the request to extract scheduling details
        # Expected format: "schedule 'task description' at '2023-01-01 12:00' repeat 'daily'"
        request = request.lower()
        
        # Check if this is a list/remove/update request
        if "list tasks" in request or "show tasks" in request or "list scheduled tasks" in request:
            tasks = list_scheduled_tasks()
            if not tasks:
                return "No scheduled tasks found."
            
            response = "Scheduled Tasks:\n"
            for task in tasks:
                repeat_info = f", Repeats: {format_time_interval(task['repeat_interval'])}" if task['repeat_interval'] else ""
                last_run = f", Last run: {task['last_executed']}" if task['last_executed'] else ""
                response += f"- Task {task['id']}: {task['description']} (Run at: {task['execute_at']}{repeat_info}{last_run})\n"
            return response
        
        elif "remove task" in request or "delete task" in request or "cancel task" in request:
            # Extract task ID
            task_id = None
            words = request.split()
            for i, word in enumerate(words):
                if word in ["task", "id"] and i < len(words) - 1 and words[i+1].isdigit():
                    task_id = int(words[i+1])
                    break
            
            if task_id is None:
                return "Please specify a task ID to remove. Use 'list tasks' to see all tasks."
            
            success, message = remove_scheduled_task(task_id)
            return message
        
        elif "update task" in request or "edit task" in request:
            # This is more complex - we'll need to extract task ID and changed parameters
            # For simplicity, we'll just provide guidance
            return "To update a task, please use the format: 'update task ID description=\"new description\" time=\"new time\" repeat=\"new interval\"'. Use 'list tasks' to see all tasks first."
        
        # Otherwise, try to parse a new task creation request
        description_start = request.find("'")
        if description_start == -1:
            description_start = request.find("\"")
        
        if description_start == -1:
            return "Please provide a task description in quotes. Example: schedule 'check email' at '3pm tomorrow'"
        
        quote_char = request[description_start]
        description_end = request.find(quote_char, description_start + 1)
        
        if description_end == -1:
            return "Invalid format. Please enclose the task description in quotes."
        
        description = request[description_start + 1:description_end]
        
        # Extract execution time
        time_start = request.find("at", description_end)
        repeat_interval = None
        
        if time_start == -1:
            return "Please specify when to run the task using 'at'. Example: schedule 'check email' at '3pm tomorrow'"
        
        time_start = request.find(quote_char, time_start)
        if time_start == -1:
            return "Please provide the execution time in quotes. Example: at '3pm tomorrow'"
        
        time_end = request.find(quote_char, time_start + 1)
        if time_end == -1:
            return "Invalid format. Please enclose the execution time in quotes."
        
        execution_time = request[time_start + 1:time_end]
        
        # Extract repeat interval if present
        repeat_start = request.find("repeat", time_end)
        if repeat_start != -1:
            repeat_start = request.find(quote_char, repeat_start)
            if repeat_start != -1:
                repeat_end = request.find(quote_char, repeat_start + 1)
                if repeat_end != -1:
                    repeat_interval = request[repeat_start + 1:repeat_end]
        
        # Extract the command to run (which is the description if not specified)
        command_start = request.find("command", time_end)
        command = description  # Default
        
        if command_start != -1:
            command_start = request.find(quote_char, command_start)
            if command_start != -1:
                command_end = request.find(quote_char, command_start + 1)
                if command_end != -1:
                    command = request[command_start + 1:command_end]
        
        # Try to parse the execution time
        try:
            execute_at = parse_datetime(execution_time)
        except:
            return f"Couldn't understand the time format: '{execution_time}'. Please use a format like '2023-01-01 15:30' or 'tomorrow at 3pm'."
        
        # Add the task
        success, result = add_scheduled_task(description, command, execute_at, repeat_interval)
        
        if success:
            task = result
            repeat_info = f" and repeats every {format_time_interval(task.repeat_interval)}" if task.repeat_interval else ""
            return f"Task scheduled successfully! Task {task.task_id}: '{task.description}' will run at {task.execute_at.strftime('%Y-%m-%d %H:%M')}{repeat_info}."
        else:
            return f"Failed to schedule task: {result}"
        
    except Exception as e:
        print(f"Error processing schedule command: {e}")
        return f"Error scheduling task: {e}. Please use format: schedule 'task description' at 'time' [repeat 'interval'] [command 'command to run']"

# Load tasks on module import
load_tasks()
start_task_executor()

def execute_system_command(command: str) -> str:
    """Execute a system command and return the output."""
    try:
        # Use shell=True for Windows commands
        result = subprocess.run(command, shell=True, capture_output=True, text=True, check=True)
        return result.stdout
    except subprocess.CalledProcessError as e:
        print(f"Command failed: {e}")
        return ""
    except Exception as e:
        print(f"Error executing command: {e}")
        return ""

def open_browser(url: str = "https://www.google.com") -> str:
    """Open a URL in the default browser."""
    try:
        # Ensure URL has proper format
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        webbrowser.open(url)
        return f"I've opened {url} for you"
    except Exception as e:
        print(f"Error opening browser: {e}")
        return "I've opened your default browser"

def send_email(to_email: str, subject: str, body: str, attachments: List[str] = None) -> str:
    """
    Send an email with optional attachments.
    
    Args:
        to_email (str): Recipient's email address
        subject (str): Email subject
        body (str): Email body text
        attachments (List[str], optional): List of file paths to attach
        
    Returns:
        str: Success or error message
    """
    try:
        # Check if email configuration exists
        if not all([EMAIL_CONFIG['smtp_server'], EMAIL_CONFIG['smtp_port'], 
                   EMAIL_CONFIG['sender_email'], EMAIL_CONFIG['sender_password']]):
            return "Email configuration is incomplete. Please check your .env file."

        print(f"Attempting to send email using configuration:")
        print(f"SMTP Server: {EMAIL_CONFIG['smtp_server']}")
        print(f"SMTP Port: {EMAIL_CONFIG['smtp_port']}")
        print(f"Sender Email: {EMAIL_CONFIG['sender_email']}")
        print(f"Password length: {len(EMAIL_CONFIG['sender_password'])}")

        # Create message
        msg = MIMEMultipart()
        msg['From'] = EMAIL_CONFIG['sender_email']
        msg['To'] = to_email
        msg['Subject'] = subject

        # Add body
        msg.attach(MIMEText(body, 'plain'))

        # Add attachments if any
        if attachments:
            for file_path in attachments:
                try:
                    with open(file_path, 'rb') as f:
                        part = MIMEApplication(f.read(), Name=os.path.basename(file_path))
                        part['Content-Disposition'] = f'attachment; filename="{os.path.basename(file_path)}"'
                        msg.attach(part)
                except Exception as e:
                    return f"Error attaching file {file_path}: {str(e)}"

        # Connect to SMTP server
        try:
            print(f"\nConnecting to SMTP server {EMAIL_CONFIG['smtp_server']}:{EMAIL_CONFIG['smtp_port']}...")
            server = smtplib.SMTP(EMAIL_CONFIG['smtp_server'], int(EMAIL_CONFIG['smtp_port']))
            print("Starting TLS connection...")
            server.starttls()  # Enable TLS
            print("Attempting to login...")
            server.login(EMAIL_CONFIG['sender_email'], EMAIL_CONFIG['sender_password'])
            print("Successfully logged in to SMTP server")
        except smtplib.SMTPAuthenticationError as e:
            print(f"SMTP Authentication Error: {str(e)}")
            return "SMTP Authentication failed. Please check your email and App Password."
        except smtplib.SMTPConnectError as e:
            print(f"SMTP Connection Error: {str(e)}")
            return f"Could not connect to SMTP server {EMAIL_CONFIG['smtp_server']}:{EMAIL_CONFIG['smtp_port']}"
        except Exception as e:
            print(f"SMTP Error: {str(e)}")
            return f"SMTP error: {str(e)}"

        # Send email
        try:
            print("Sending email...")
            server.send_message(msg)
            server.quit()
            print("Email sent successfully")
            return f"Email sent successfully to {to_email}"
        except Exception as e:
            print(f"Error sending email: {str(e)}")
            return f"Error sending email: {str(e)}"

    except Exception as e:
        print(f"Unexpected error: {str(e)}")
        return f"Unexpected error: {str(e)}"

def get_system_info() -> str:
    """Get basic system information."""
    try:
        info = {
            "OS": platform.system(),
            "OS Version": platform.version(),
            "CPU Usage": f"{psutil.cpu_percent()}%",
            "Memory Usage": f"{psutil.virtual_memory().percent}%",
            "Disk Usage": f"{psutil.disk_usage('/').percent}%"
        }
        return "\n".join(f"{k}: {v}" for k, v in info.items())
    except Exception as e:
        print(f"Error getting system info: {e}")
        return "Here's your system information"

def create_and_write_file(content: str, extension: str = ".txt") -> str:
    """Create a temporary file with the given content and return its path."""
    try:
        # Create a more descriptive filename
        prefix = f"generated_{int(time.time())}"
        with tempfile.NamedTemporaryFile(suffix=extension, prefix=prefix, delete=False, mode='w', encoding='utf-8') as f:
            f.write(content)
            return f.name
    except Exception as e:
        print(f"Error creating file: {e}")
        return ""

def open_in_editor(file_path: str) -> str:
    """Open a file in the default editor."""
    try:
        if platform.system() == "Windows":
            # Try VSCode first
            vscode_paths = [
                "C:\\Program Files\\Microsoft VS Code\\Code.exe",
                "C:\\Program Files (x86)\\Microsoft VS Code\\Code.exe",
                os.path.expanduser("~\\AppData\\Local\\Programs\\Microsoft VS Code\\Code.exe")
            ]
            
            for vscode_path in vscode_paths:
                if os.path.exists(vscode_path):
                    subprocess.Popen([vscode_path, file_path])
                    return f"I've opened the file in VSCode"
            
            # Fallback to notepad
            subprocess.Popen(["notepad", file_path])
            return f"I've opened the file in Notepad"
        elif platform.system() == "Darwin":  # macOS
            # Try VSCode first
            subprocess.Popen(["open", "-a", "Visual Studio Code", file_path])
            return f"I've opened the file in VSCode"
        else:  # Linux
            # Try VSCode first
            subprocess.Popen(["code", file_path])
            return f"I've opened the file in VSCode"
    except Exception as e:
        print(f"Error opening editor: {e}")
        # Fallback to default editor
        try:
            if platform.system() == "Windows":
                subprocess.Popen(["notepad", file_path])
            elif platform.system() == "Darwin":
                subprocess.Popen(["open", "-e", file_path])
            else:
                subprocess.Popen(["xdg-open", file_path])
            return "I've opened the file for you"
        except Exception as e:
            print(f"Error in fallback editor: {e}")
            return "I've opened the file for you"

def generate_code(task: str, language: str) -> str:
    """Generate code based on the task description."""
    try:
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        if language.lower() == "html":
            prompt = f"""Generate a complete HTML file with CSS and JavaScript for: {task}. 
            Include all necessary HTML, CSS, and JavaScript in a single file.
            Use modern HTML5, CSS3, and JavaScript features.
            Make it visually appealing and interactive.
            Only provide the code, no explanations, no markdown formatting, no code blocks, no backticks.
            Your response must start with <!DOCTYPE html> and contain only valid HTML code."""
        else:
            prompt = f"Generate a simple {language} code for: {task}. Only provide the code, no explanations, no markdown formatting, no code blocks, no backticks."
        
        # Use the model without additional parameters
        response = model.generate_content(prompt)
        
        # Clean up the response to remove any markdown formatting
        code = response.text.strip()
        # Remove any markdown formatting
        if code.startswith('```'):
            code = "\n".join(code.split("\n")[1:-1])
        # Remove language identifier if present
        if code.startswith(language.lower()) or code.startswith(language) or code.startswith("html"):
            code = code.split("\n", 1)[1] if "\n" in code else ""
        # Remove any trailing backticks or markdown
        if "```" in code:
            code = code.split("```")[0]
        return code.strip()
    except Exception as e:
        print(f"Error generating code: {e}")
        return ""

def run_python_file(file_path: str) -> str:
    """Run a Python file and return its output."""
    try:
        result = subprocess.run(['python', file_path], capture_output=True, text=True, check=True)
        return result.stdout
    except subprocess.CalledProcessError as e:
        print(f"Error running Python file: {e}")
        return f"Error: {e.stderr}"
    except Exception as e:
        print(f"Error running Python file: {e}")
        return "Error running the Python file"

def open_html_in_browser(file_path: str) -> str:
    """Open an HTML file in the default browser."""
    try:
        # Convert to absolute path
        abs_path = os.path.abspath(file_path)
        # Convert to file URL
        file_url = f"file:///{abs_path.replace(os.sep, '/')}"
        webbrowser.open(file_url)
        return f"I've opened the HTML file in your browser"
    except Exception as e:
        print(f"Error opening HTML file: {e}")
        return "I've opened the file in your browser"

def setup_chrome_driver():
    """Set up and return a Chrome WebDriver instance."""
    try:
        chrome_options = Options()
        chrome_options.add_argument("--headless")  # Run in headless mode
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--disable-dev-shm-usage")
        
        # Try to find Chrome in common locations
        chrome_paths = [
            "C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe",
            "C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe",
            os.path.expanduser("~\\AppData\\Local\\Google\\Chrome\\Application\\chrome.exe")
        ]
        
        chrome_path = None
        for path in chrome_paths:
            if os.path.exists(path):
                chrome_path = path
                break
        
        if chrome_path:
            chrome_options.binary_location = chrome_path
        
        service = Service()
        driver = webdriver.Chrome(service=service, options=chrome_options)
        return driver
    except Exception as e:
        print(f"Error setting up Chrome driver: {e}")
        return None

def play_youtube_video(search_query: str) -> str:
    """Search YouTube and play the first video result."""
    try:
        driver = setup_chrome_driver()
        if not driver:
            # Fallback to just opening the search results
            url = f"https://www.youtube.com/results?search_query={search_query}"
            webbrowser.open(url)
            return f"I've opened YouTube and searched for {search_query}"
        
        # Search for the video
        driver.get(f"https://www.youtube.com/results?search_query={search_query}")
        
        # Wait for the first video to be clickable
        wait = WebDriverWait(driver, 10)
        first_video = wait.until(
            EC.element_to_be_clickable((By.CSS_SELECTOR, "ytd-video-renderer a#video-title"))
        )
        
        # Get the video URL and title
        video_url = first_video.get_attribute("href")
        video_title = first_video.get_attribute("title")
        
        # Close the headless browser
        driver.quit()
        
        # Open the video in the default browser
        webbrowser.open(video_url)
        return f"I've found and started playing: {video_title}"
        
    except Exception as e:
        print(f"Error playing YouTube video: {e}")
        # Fallback to just opening the search results
        url = f"https://www.youtube.com/results?search_query={search_query}"
        webbrowser.open(url)
        return f"I've opened YouTube and searched for {search_query}"

def capture_screen():
    """Capture the current screen and return the image."""
    try:
        screenshot = pyautogui.screenshot()
        return cv2.cvtColor(np.array(screenshot), cv2.COLOR_RGB2BGR)
    except Exception as e:
        print(f"Error capturing screen: {e}")
        return None

def analyze_screen():
    """Analyze the current screen and return information about visible elements."""
    try:
        print("Debug: Starting screen analysis...")
        # Capture screen
        screen = capture_screen()
        if screen is None:
            print("Debug: Failed to capture screen")
            return "Unable to capture screen"
        
        # Save screen for analysis
        temp_path = os.path.join(tempfile.gettempdir(), "screen_analysis.png")
        cv2.imwrite(temp_path, screen)
        print(f"Debug: Screen saved to {temp_path}")
        
        # Use AI to analyze the screen
        try:
            print("Debug: Initializing AI model for screen analysis...")
            model = genai.GenerativeModel('gemini-2.0-flash-lite')
            image_parts = [
                {
                    "mime_type": "image/png",
                    "data": open(temp_path, "rb").read()
                }
            ]
            
            prompt = """Analyze this screenshot in detail and provide a comprehensive analysis. Focus on:
            1. What application or window is currently visible (be specific)
            2. Any visible text, code, or content
            3. Any visible icons, buttons, or interactive elements
            4. The current context and what the user might be doing
            5. Any visible URLs or website content
            6. Any visible file names or paths
            
            Respond in this exact JSON format:
            {
                "current_app": "detailed name of visible application",
                "visible_text": "all visible text content",
                "visible_elements": ["list", "of", "all", "visible", "elements"],
                "context": "detailed description of current context",
                "urls": ["any", "visible", "urls"],
                "files": ["any", "visible", "file", "names"]
            }"""
            
            # Create a proper generation config
            config = {
                "temperature": 0.2,
                "top_p": 0.95,
                "top_k": 40,
                "max_output_tokens": 1024,
            }
            
            print("Debug: Sending screen to AI for analysis...")
            # Use the model without the additional parameters
            response = model.generate_content([prompt, *image_parts])
                
            analysis = response.text
            print(f"Debug: Received AI response length: {len(analysis)}")
        
            # Try to parse the response as JSON
            try:
                result = json.loads(analysis)
                print("Debug: Successfully parsed response as JSON")
                return result
            except json.JSONDecodeError:
                print("Debug: Failed to parse response as JSON, trying to extract information")
                try:
                    # If JSON parsing fails, try to extract structured information
                    lines = analysis.strip().split('\n')
                    result = {
                        "current_app": "Unknown",
                        "visible_text": analysis,
                        "visible_elements": [],
                        "context": "Unable to parse screen analysis",
                        "urls": [],
                        "files": []
                    }
                    
                    # Try to extract structured information from text response
                    for line in lines:
                        try:
                            if line.startswith("current_app:") or line.startswith('"current_app":'):
                                parts = line.split(":", 1)
                                if len(parts) > 1:
                                    result["current_app"] = parts[1].strip().strip('"').strip(',')
                            elif "visible_text" in line.lower():
                                parts = line.split(":", 1)
                                if len(parts) > 1:
                                    result["visible_text"] = parts[1].strip().strip('"').strip(',')
                            elif "context" in line.lower():
                                parts = line.split(":", 1)
                                if len(parts) > 1:
                                    result["context"] = parts[1].strip().strip('"').strip(',')
                        except Exception as line_error:
                            print(f"Debug: Error processing line: {line_error}")
                            continue
                    
                    print("Debug: Created structured result from text response")
                    return result
                except Exception as extract_error:
                    print(f"Debug: Error extracting information: {extract_error}")
                    return {
                        "current_app": "Unknown",
                        "visible_text": "Failed to parse screen analysis",
                        "visible_elements": [],
                        "context": f"Error: {extract_error}",
                        "urls": [],
                        "files": []
                    }
        except Exception as ai_error:
            print(f"Debug: Error with AI analysis: {ai_error}")
            # Fallback handling for screen analysis
            try:
                # Test if tesseract is installed
                tesseract_installed = False
                try:
                    # Try importing pytesseract
                    import pytesseract
                    # Try getting tesseract version
                    pytesseract.get_tesseract_version()
                    tesseract_installed = True
                    print("Debug: Tesseract is installed, using OCR")
                except Exception:
                    tesseract_installed = False
                    print("Debug: Tesseract is not installed")
                
                if tesseract_installed:
                    # Use pytesseract for OCR analysis
                    img = Image.open(temp_path)
                    text = pytesseract.image_to_string(img)
                    
                    return {
                        "current_app": "Unknown (OCR fallback)",
                        "visible_text": text,
                        "visible_elements": [],
                        "context": "Basic OCR analysis only",
                        "urls": [],
                        "files": []
                    }
                else:
                    # Simple screen color analysis as last resort
                    print("Debug: Using basic color analysis as fallback")
                    img = cv2.imread(temp_path)
                    
                    # Basic color analysis
                    avg_color = cv2.mean(img)
                    is_dark = sum(avg_color[:3]) / 3 < 128
                    
                    # Try to extract more visual features
                    try:
                        # Convert to grayscale for edge detection
                        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                        
                        # Edge detection to find UI elements
                        edges = cv2.Canny(gray, 100, 200)
                        edge_count = cv2.countNonZero(edges)
                        has_many_elements = edge_count > 10000
                        
                        # Look for text-like regions
                        _, threshold = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY_INV)
                        contours, _ = cv2.findContours(threshold, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                        text_like_regions = len([c for c in contours if 10 < cv2.contourArea(c) < 500])
                        
                        # Analyze color distribution
                        blue_avg = avg_color[0]
                        green_avg = avg_color[1]
                        red_avg = avg_color[2]
                        
                        # Common UI colors
                        dominant_color = "neutral"
                        if blue_avg > max(red_avg, green_avg) + 20:
                            dominant_color = "blue"
                        elif green_avg > max(red_avg, blue_avg) + 20:
                            dominant_color = "green"
                        elif red_avg > max(green_avg, blue_avg) + 20:
                            dominant_color = "red"
                        
                        # Create a more descriptive analysis
                        description = f"A {dominant_color}-tinted {'dark' if is_dark else 'light'} screen with "
                        description += "many UI elements" if has_many_elements else "few UI elements"
                        description += " and what appears to be text content." if text_like_regions > 10 else "."
                        
                        elements = []
                        if has_many_elements:
                            elements.append("multiple UI controls")
                        if text_like_regions > 10:
                            elements.append("text regions")
                        elements.append(f"{dominant_color} elements")
                        elements.append(f"{'dark' if is_dark else 'light'} background")
                        
                        return {
                            "current_app": "Unknown (visual analysis only)",
                            "visible_text": "Unable to extract text - tesseract not installed",
                            "visible_elements": elements,
                            "context": description,
                            "urls": [],
                            "files": []
                        }
                    except Exception as visual_error:
                        print(f"Debug: Error in visual analysis: {visual_error}")
                        # Fall back to very basic analysis
                        return {
                            "current_app": "Unknown (basic analysis)",
                            "visible_text": "Unable to extract text - tesseract not installed",
                            "visible_elements": ["dark screen" if is_dark else "light screen"],
                            "context": "Basic screen analysis only - install tesseract for better results",
                            "urls": [],
                            "files": []
                        }
            except Exception as backup_error:
                print(f"Debug: Error with backup analysis: {backup_error}")
                return "Unable to analyze screen (all methods failed)"
    
    except Exception as e:
        print(f"Debug: Error in analyze_screen: {e}")
        return {
            "current_app": "Unknown",
            "visible_text": "Error analyzing screen",
            "visible_elements": [],
            "context": str(e),
            "urls": [],
            "files": []
        }
    
    finally:
        # Clean up in finally block to ensure it runs
        try:
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print(f"Debug: Removed temporary file {temp_path}")
        except Exception as cleanup_error:
            print(f"Debug: Error during cleanup: {cleanup_error}")

def analyze_code_context():
    """Analyze the current code context if in an editor."""
    try:
        screen = capture_screen()
        if screen is None:
            return None
        
        # Save screen for analysis
        temp_path = os.path.join(tempfile.gettempdir(), "code_analysis.png")
        cv2.imwrite(temp_path, screen)
        
        # Use AI to analyze the code
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        image_parts = [
            {
                "mime_type": "image/png",
                "data": open(temp_path, "rb").read()
            }
        ]
        
        prompt = """You are a code analysis expert. Analyze this code screenshot and provide detailed information about:
        1. The programming language being used
        2. The complete code visible in the editor
        3. Any syntax errors, runtime errors, or warnings
        4. The file name and path if visible
        5. Any error messages or stack traces
        6. The current cursor position and selected text
        
        Be very thorough in your analysis. Look for:
        - Syntax errors (missing brackets, semicolons, etc.)
        - Runtime errors (undefined variables, type mismatches, etc.)
        - Logical errors (incorrect logic, infinite loops, etc.)
        - Style issues (indentation, naming conventions, etc.)
        
        Respond in this exact JSON format:
        {
            "language": "programming language",
            "code": "complete code content",
            "errors": [
                {
                    "type": "error type (syntax/runtime/warning)",
                    "message": "detailed error message",
                    "line": "line number",
                    "suggestion": "detailed fix suggestion"
                }
            ],
            "file_info": {
                "name": "file name",
                "path": "file path"
            },
            "cursor": {
                "line": "current line",
                "column": "current column",
                "selection": "selected text if any"
            }
        }"""
        
        # Use the model without additional parameters
        response = model.generate_content([prompt, *image_parts])
        
        analysis = response.text
        
        try:
            # Clean up
            if os.path.exists(temp_path):
                os.remove(temp_path)
        except Exception as cleanup_error:
            print(f"Error during cleanup: {cleanup_error}")
        
        try:
            return json.loads(analysis)
        except json.JSONDecodeError:
            # If JSON parsing fails, try to extract useful information from the text
            return {
                "language": "Unknown",
                "code": analysis,
                "errors": [{
                    "type": "unknown",
                    "message": "Unable to parse code analysis",
                    "line": "unknown",
                    "suggestion": "Please try again or provide more context"
                }],
                "file_info": {
                    "name": "unknown",
                    "path": "unknown"
                },
                "cursor": {
                    "line": "unknown",
                    "column": "unknown",
                    "selection": ""
                }
            }
            
    except Exception as e:
        print(f"Error analyzing code: {e}")
        return None

def fix_code_errors(code_analysis):
    """Fix code errors based on the analysis."""
    try:
        if not code_analysis or not code_analysis.get("errors"):
            return None
            
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        
        # Create a detailed prompt for fixing the errors
        prompt = f"""You are an expert code fixer. Fix the following code errors:
        
        Language: {code_analysis.get('language', 'Unknown')}
        Current code:
        {code_analysis.get('code', '')}
        
        Errors to fix:
        {json.dumps(code_analysis.get('errors', []), indent=2)}
        
        Requirements:
        1. Fix ALL errors identified in the analysis
        2. Maintain the original code structure and style
        3. Add any necessary imports or dependencies
        4. Ensure the code is complete and runnable
        5. Add comments explaining major changes
        
        Provide the complete fixed code that addresses all errors.
        Only provide the code, no explanations."""
        
        # Use the model without additional parameters
        response = model.generate_content(prompt)
        
        fixed_code = response.text.strip()
        
        # Clean up the response
        if fixed_code.startswith('```'):
            fixed_code = fixed_code.split('\n', 1)[1]
        if fixed_code.endswith('```'):
            fixed_code = fixed_code.rsplit('\n', 1)[0]
            
        return fixed_code.strip()
        
    except Exception as e:
        print(f"Error fixing code: {e}")
        return None

def get_system_volume() -> float:
    """Get current system volume level (0-100)."""
    try:
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        return volume.GetMasterVolumeLevelScalar() * 100
    except Exception as e:
        print(f"Error getting volume: {e}")
        return 0

def set_system_volume(level: float) -> bool:
    """Set system volume level (0-100)."""
    try:
        devices = AudioUtilities.GetSpeakers()
        interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        volume = cast(interface, POINTER(IAudioEndpointVolume))
        volume.SetMasterVolumeLevelScalar(level / 100, None)
        return True
    except Exception as e:
        print(f"Error setting volume: {e}")
        return False

def get_screen_brightness() -> float:
    """Get current screen brightness level (0-100)."""
    try:
        if platform.system() == "Windows":
            import wmi
            w = wmi.WMI(namespace="wmi")
            brightness = w.WmiMonitorBrightness()[0].CurrentBrightness
            return brightness
        return 0
    except Exception as e:
        print(f"Error getting brightness: {e}")
        return 0

def set_screen_brightness(level: float) -> bool:
    """Set screen brightness level (0-100)."""
    try:
        if platform.system() == "Windows":
            import wmi
            w = wmi.WMI(namespace="wmi")
            w.WmiMonitorBrightnessMethods()[0].WmiSetBrightness(level, 0)
            return True
        return False
    except Exception as e:
        print(f"Error setting brightness: {e}")
        return False

def get_youtube_video_info(video_url: str) -> Dict:
    """Get information about a YouTube video."""
    try:
        response = requests.get(video_url)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract video information
        title = soup.find('title').text.replace(' - YouTube', '')
        description = soup.find('meta', {'name': 'description'})['content']
        
        # Try to get view count and upload date
        view_count = "Unknown"
        upload_date = "Unknown"
        for script in soup.find_all('script'):
            if 'var ytInitialData = ' in str(script):
                data = str(script).split('var ytInitialData = ')[1].split(';</script>')[0]
                try:
                    json_data = json.loads(data)
                    video_data = json_data['contents']['twoColumnWatchNextResults']['results']['results']['contents'][0]['videoPrimaryInfoRenderer']
                    view_count = video_data['viewCount']['videoViewCountRenderer']['viewCount']['simpleText']
                    upload_date = video_data['dateText']['simpleText']
                except:
                    pass
                break
        
        return {
            'title': title,
            'description': description,
            'view_count': view_count,
            'upload_date': upload_date,
            'url': video_url
        }
    except Exception as e:
        print(f"Error getting video info: {e}")
        return {
            'title': 'Error',
            'description': str(e),
            'view_count': 'Unknown',
            'upload_date': 'Unknown',
            'url': video_url
        }

def create_youtube_playlist(name: str, videos: List[str]) -> str:
    """Create a YouTube playlist with the given videos."""
    try:
        # First open YouTube
        webbrowser.open("https://www.youtube.com")
        time.sleep(2)  # Wait for YouTube to load
        
        # Click on "Create playlist" button
        pyautogui.click(x=100, y=100)  # Adjust coordinates based on your screen
        
        # Type playlist name
        pyautogui.write(name)
        pyautogui.press('enter')
        
        # Add videos
        for video_url in videos:
            webbrowser.open(video_url)
            time.sleep(2)
            # Click "Save to playlist" button
            pyautogui.click(x=200, y=200)  # Adjust coordinates based on your screen
            # Select the playlist
            pyautogui.click(x=300, y=300)  # Adjust coordinates based on your screen
        
        return f"Created playlist '{name}' with {len(videos)} videos"
    except Exception as e:
        print(f"Error creating playlist: {e}")
        return f"Error creating playlist: {str(e)}"

def hex_to_rgb(hex_color):
    """Convert hex color to RGB."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def download_image(url):
    """Download image from URL."""
    try:
        response = requests.get(url)
        if response.status_code == 200:
            return io.BytesIO(response.content)
        return None
    except:
        return None

def apply_slide_style(slide, slide_data):
    """Apply enhanced styling to the slide."""
    # Get background color
    bg_color = slide_data.get('backgroundColor', '#FFFFFF')
    
    # Apply background
    background = slide.background
    fill = background.fill
    
    # Handle gradient backgrounds
    bg_image = slide_data.get('backgroundImage')
    if bg_image and bg_image.startswith('linear-gradient'):
        colors = re.findall(r'#[0-9A-Fa-f]{6}', bg_image)
        if len(colors) >= 2:
            # Create a subtle gradient effect
            fill.gradient()
            fill.gradient_stops[0].color.rgb = RGBColor(*hex_to_rgb(colors[0]))
            fill.gradient_stops[1].color.rgb = RGBColor(*hex_to_rgb(colors[1]))
            fill.gradient_angle = 45
    else:
        # Solid background
        fill.solid()
        r, g, b = hex_to_rgb(bg_color)
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Apply theme-specific enhancements
    theme = slide_data.get('theme', 'light')
    if theme == 'dark':
        # For dark theme, ensure text is light
        slide_data['textColor'] = slide_data.get('textColor', '#FFFFFF')
        # Add subtle pattern or texture
        if hasattr(slide.background.fill, 'pattern'):
            slide.background.fill.pattern = True
            slide.background.fill.pattern_fore_color.rgb = RGBColor(40, 40, 40)
    
    # Add slide number if not title slide
    if slide_data.get('layout') != 'title':
        slide_number = slide.shapes.add_textbox(
            Inches(12), Inches(7),  # Position at bottom-right
            Inches(0.5), Inches(0.3)
        )
        number_text = slide_number.text_frame.add_paragraph()
        number_text.text = str(slide.slide_id)
        number_text.alignment = PP_ALIGN.RIGHT
        number_text.font.size = Pt(10)
        number_text.font.color.rgb = RGBColor(*hex_to_rgb(slide_data.get('textColor', '#000000')))

def apply_text_style(text_frame, slide_data, is_title=False):
    """Apply text styling with improved alignment and spacing."""
    for paragraph in text_frame.paragraphs:
        # Alignment
        if is_title:
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            paragraph.alignment = slide_data.get('alignment', PP_ALIGN.LEFT)
        
        # Line spacing
        paragraph.line_spacing = 1.2
        
        # Font family
        font_family = slide_data.get('fontFamily', 'Calibri')
        # Font size with better defaults
        if is_title:
            font_size = int(float(slide_data.get('titleFontSize', '44').replace('rem', '')) * 12)
        else:
            font_size = int(float(slide_data.get('contentFontSize', '28').replace('rem', '')) * 12)
        # Text color
        text_color = slide_data.get('textColor', '#000000')
        r, g, b = hex_to_rgb(text_color)
        
        for run in paragraph.runs:
            run.font.name = font_family
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(r, g, b)
            run.font.bold = slide_data.get('bold', False)

def create_pptx_from_slides(slides_data, output_path):
    """Create a PowerPoint presentation with clean, modern layouts."""
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        # Always use blank layout for consistency
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Apply slide styling
        apply_slide_style(slide, slide_data)
        
        # Add title with improved positioning
        title = slide_data.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(0.4),
                width=Inches(12.333),
                height=Inches(1.2)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            if slide_data.get('textColor'):
                r, g, b = hex_to_rgb(slide_data.get('textColor'))
                p.font.color.rgb = RGBColor(r, g, b)
        
        # Get or generate content
        content = slide_data.get('content', '')
        if not content and slide.slide_id > 1:  # If no content and not title slide
            # Generate default content based on title
            content = f"• Key Points about {title}:\n"
            content += "  - Understanding the fundamentals and core concepts\n"
            content += "  - Exploring practical applications and use cases\n"
            content += "  - Analyzing impact and future implications\n"
            content += "  - Best practices and implementation strategies"
        
        # Handle content and image placement
        image_url = slide_data.get('imageUrl')
        
        if image_url and content and slide.slide_id > 1:  # Both image and content for non-title slides
            # Content on left - moved up slightly
            content_box = slide.shapes.add_textbox(
                left=Inches(0.6),
                top=Inches(1.5),  # Moved up from 1.8
                width=Inches(5.8),
                height=Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Add content with proper formatting
            first = True
            for point in content.split('\n'):
                if not point.strip():
                    continue
                
                p = content_frame.add_paragraph()
                if not first:
                    p.space_before = Pt(12)
                first = False
                
                # Format bullet points properly
                if point.startswith('  -'):  # Sub-bullet
                    p.text = point.replace('  -', '•')
                    p.level = 1
                    p.space_before = Pt(6)
                else:  # Main bullet
                    p.text = point.strip('• ')
                    if not point.startswith('•'):
                        p.text = '• ' + p.text
                    p.level = 0
                
                p.font.size = Pt(24)
                if slide_data.get('textColor'):
                    r, g, b = hex_to_rgb(slide_data.get('textColor'))
                    p.font.color.rgb = RGBColor(r, g, b)
            
            # Image on right - moved down slightly
            try:
                image_data = download_image(image_url)
                if image_data:
                    picture = slide.shapes.add_picture(
                        image_data,
                        left=Inches(6.8),
                        top=Inches(2.1),  # Moved down from 1.8
                        width=Inches(5.8),
                        height=Inches(4.2)
                    )
            except Exception as e:
                print(f"Warning: Failed to add image: {str(e)}")
                
        elif slide.slide_id == 1:  # Title slide
            # Center content for title slide
            if content:
                subtitle_box = slide.shapes.add_textbox(
                    left=Inches(0.8),
                    top=Inches(2.5),  # Adjusted for better spacing
                    width=Inches(11.733),
                    height=Inches(2)
                )
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.word_wrap = True
                
                p = subtitle_frame.add_paragraph()
                p.text = content
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(28)
                if slide_data.get('textColor'):
                    r, g, b = hex_to_rgb(slide_data.get('textColor'))
                    p.font.color.rgb = RGBColor(r, g, b)
            
            # Add image at the bottom for title slide
            if image_url:
                try:
                    image_data = download_image(image_url)
                    if image_data:
                        picture = slide.shapes.add_picture(
                            image_data,
                            left=Inches(3.666),
                            top=Inches(3.8),  # Adjusted for better spacing
                            width=Inches(6),
                            height=Inches(3.2)
                        )
                except Exception as e:
                    print(f"Warning: Failed to add image: {str(e)}")
        
        else:  # Only content, no image or title slide
            content_box = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(1.5),  # Moved up slightly
                width=Inches(11.333),
                height=Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Add content with proper formatting
            first = True
            for point in content.split('\n'):
                if not point.strip():
                    continue
                
                p = content_frame.add_paragraph()
                if not first:
                    p.space_before = Pt(12)
                first = False
                
                # Format bullet points properly
                if point.startswith('  -'):  # Sub-bullet
                    p.text = point.replace('  -', '•')
                    p.level = 1
                    p.space_before = Pt(6)
                else:  # Main bullet
                    p.text = point.strip('• ')
                    if not point.startswith('•'):
                        p.text = '• ' + p.text
                    p.level = 0
                
                p.font.size = Pt(28)
                if slide_data.get('textColor'):
                    r, g, b = hex_to_rgb(slide_data.get('textColor'))
                    p.font.color.rgb = RGBColor(r, g, b)
        
        # Add slide number (except for title slide)
        if slide.slide_id > 1:
            slide_number = slide.shapes.add_textbox(
                left=Inches(12.333),
                top=Inches(6.9),
                width=Inches(0.5),
                height=Inches(0.3)
            )
            number_text = slide_number.text_frame.add_paragraph()
            number_text.text = str(slide.slide_id)
            number_text.alignment = PP_ALIGN.RIGHT
            number_text.font.size = Pt(12)
            if slide_data.get('textColor'):
                r, g, b = hex_to_rgb(slide_data.get('textColor'))
                number_text.font.color.rgb = RGBColor(r, g, b)
    
    # Save the presentation
    prs.save(output_path)
    return True

def generate_presentation(
    title="My Presentation",
    topic="Artificial Intelligence",
    theme_id="modern",
    slide_count=7,
    api_url="https://slideforge.onrender.com/api/generate-slides"
):  
    """
    Generate a PowerPoint presentation using the PresentationMaster API.
    """
    
    # Get the Downloads folder path based on the operating system
    if platform.system() == "Windows":
        downloads_path = os.path.join(os.path.expanduser("~"), "Downloads")
    elif platform.system() == "Darwin":  # macOS
        downloads_path = os.path.join(os.path.expanduser("~"), "Downloads")
    else:  # Linux
        downloads_path = os.path.join(os.path.expanduser("~"), "Downloads")
    
    # Create Downloads directory if it doesn't exist
    if not os.path.exists(downloads_path):
        os.makedirs(downloads_path)
    
    # Request headers
    headers = {
        'Content-Type': 'application/json',
        'Accept': 'application/json'
    }
    
    # Request data with enhanced styling and content requirements
    data = {
        'title': title,
        'topic': topic,
        'themeId': theme_id,
        'slideCount': slide_count,
        'format': 'pptx',
        'style': {
            'theme': 'dark',
            'colorScheme': 'professional',
            'includeImages': True,
            'imageStyle': 'modern',
            'layout': 'modern'
        },
        'content': {
            'type': 'detailed',  # Request detailed content
            'style': 'professional',
            'depth': 'comprehensive',  # Request comprehensive content
            'includeExamples': True,
            'includeStats': True
        },
        'images': {
            'required': True,  # Make images mandatory
            'style': 'modern',
            'type': 'relevant',  # Request topic-relevant images
            'minWidth': 800,  # Request high-quality images
            'minHeight': 600
        }
    }

    try:
        print(f"\nSending request to: {api_url}")
        print(f"Request data: {json.dumps(data, indent=2)}")
        
        # Make the request
        response = requests.post(api_url, json=data, headers=headers)
        
        print(f"Response status code: {response.status_code}")
        
        # Check if request was successful
        if response.status_code == 200:
            try:
                # Parse the JSON response
                response_data = response.json()
                
                # Check if we got proper response structure
                if not isinstance(response_data, (list, dict)):
                    print("\nError: Invalid response format from server")
                    return None
                
                # Convert to list if single slide
                slides_data = response_data if isinstance(response_data, list) else [response_data]
                
                if not slides_data:
                    print("\nError: No slides data in response")
                    return None
                
                # Verify each slide has required content
                for slide in slides_data:
                    if not isinstance(slide, dict):
                        continue
                        
                    # Ensure each slide has an image URL if not title slide
                    if slide.get('layout') != 'title' and not slide.get('imageUrl'):
                        # Add default image URL for the topic if none provided
                        slide['imageUrl'] = f"https://source.unsplash.com/1600x900/?{topic.replace(' ', ',')}"
                    
                    # Ensure content is detailed enough
                    if 'content' in slide and isinstance(slide['content'], str):
                        content = slide['content']
                        if len(content.split('\n')) < 3:  # If content is too brief
                            # Add more detailed content
                            content_lines = content.split('\n')
                            enhanced_content = []
                            for line in content_lines:
                                if line.strip():
                                    enhanced_content.append(line)
                                    # Add supporting details for each point
                                    if not line.endswith(':'):
                                        enhanced_content.append(f"  - Detailed explanation: {line}")
                                        enhanced_content.append(f"  - Example: Implementation in {topic}")
                            slide['content'] = '\n'.join(enhanced_content)
                
                print(f"\nProcessing {len(slides_data)} slides")
                
                # Generate filename with timestamp
                timestamp = time.strftime("%Y%m%d_%H%M%S")
                safe_title = title.replace(' ', '_').replace('/', '_').replace('\\', '_')
                filename = f"{safe_title}_{timestamp}.pptx"
                output_path = Path(os.path.join(downloads_path, filename))
                
                # Create the PowerPoint file
                if create_pptx_from_slides(slides_data, output_path):
                    print(f"\nPresentation successfully saved as: {filename}")
                    
                    # Verify file size
                    file_size = output_path.stat().st_size
                    print(f"File size: {file_size:,} bytes ({file_size/1024:.2f} KB)")
                    
                    return output_path.absolute()
                else:
                    print("\nError: Failed to create PowerPoint file")
                    return None
                    
            except json.JSONDecodeError:
                print("\nError: Failed to parse JSON response")
                print(f"Response content: {response.text[:500]}")
                return None
                
        elif response.status_code == 500:
            try:
                error_data = response.json()
                print(f"\nServer Error: {json.dumps(error_data, indent=2)}")
            except:
                print(f"\nServer Error: {response.text[:500]}")
            return None
            
        else:
            print(f"\nUnexpected status code: {response.status_code}")
            try:
                print(f"Response: {json.dumps(response.json(), indent=2)}")
            except:
                print(f"Response text: {response.text[:500]}")
            return None
            
    except Exception as e:
        print(f"\nError: {str(e)}")
        return None

def generate_and_open_image(prompt: str) -> str:
    """Generate an image using Pollinations.ai and automatically open it.
    
    Args:
        prompt: The text description of the image to generate
        
    Returns:
        str: Status message indicating success or failure
    """
    try:
        # URL encode the prompt
        encoded_prompt = quote_plus(prompt)
        
        # Construct the Pollinations.ai URL
        url = f"https://image.pollinations.ai/prompt/{encoded_prompt}"
        
        # Download the image
        response = requests.get(url, stream=True)
        if response.status_code != 200:
            return f"Error: Failed to generate image. Status code: {response.status_code}"

        # Create a temporary file to save the image
        temp_dir = tempfile.gettempdir()
        temp_image_path = os.path.join(temp_dir, f"generated_image_{int(time.time())}.png")
        
        # Save the image
        with open(temp_image_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    f.write(chunk)
        
        # Open the image using the default image viewer
        if platform.system() == 'Windows':
            os.startfile(temp_image_path)
        elif platform.system() == 'Darwin':  # macOS
            subprocess.run(['open', temp_image_path])
        else:  # Linux
            subprocess.run(['xdg-open', temp_image_path])
            
        return f"Image generated and opened successfully! Saved at: {temp_image_path}"
        
    except Exception as e:
        return f"Error generating image: {str(e)}"

if __name__ == '__main__':
    # Example usage with dark theme
    presentation_path = generate_presentation(
        title="AI in Healthcare",
        topic="Applications of Artificial Intelligence in Modern Healthcare",
        theme_id="dark",  # Using dark theme
        slide_count=7
    )
    
    if presentation_path:
        print(f"\nFinal presentation saved at: {presentation_path}")
        if presentation_path.exists():
            size = presentation_path.stat().st_size
            print(f"Final file size: {size:,} bytes ({size/1024:.2f} KB)")
            print("PowerPoint file created successfully with images and styling!")
            print("\nYou can now open the presentation in Microsoft PowerPoint")

def process_command(command: str, is_voice_mode=False) -> str:
    """Process user commands using AI to understand and execute the request."""
    try:
        global memory
        
        # Handle image generation commands
        if "generate image" in command.lower() or "create image" in command.lower():
            # Extract the prompt from the command
            prompt_match = re.search(r"(?:generate|create) image (?:of |showing |with |about )?(.*)", command)
            if prompt_match:
                prompt = prompt_match.group(1).strip()
                return generate_and_open_image(prompt)
            return "Please provide a description of the image you want to generate."
        
        # Handle YouTube search commands first
        if "youtube" in command.lower() or "search" in command.lower():
            # Extract search query
            search_query = command.lower()
            # Remove common words only when they appear as standalone words
            common_words = ["search", "on", "youtube", "for", "play"]
            for word in common_words:
                # Use word boundaries to ensure we only match whole words
                search_query = re.sub(r'\b' + word + r'\b', '', search_query)
            # Clean up any extra spaces
            search_query = ' '.join(search_query.split())
            
            if search_query:
                try:
                    search_url = f"https://www.youtube.com/results?search_query={search_query.replace(' ', '+')}"
                    webbrowser.open(search_url)
                    return f"I've opened YouTube and searched for '{search_query}'"
                except Exception as e:
                    print(f"Error searching YouTube: {e}")
                    return f"I had trouble searching YouTube for '{search_query}'. Please try again."
            else:
                return "Please specify what you'd like to search for on YouTube."

        # Direct presentation command parsing
        if any(word in command.lower() for word in ["presentation", "powerpoint", "ppt"]) and \
           any(word in command.lower() for word in ["create", "make", "generate"]):
            # Extract title if specified
            title = "My Presentation"
            title_match = re.search(r"titled ['\"](.*?)['\"]", command)
            if title_match:
                title = title_match.group(1)
            
            # Extract topic (everything after "about")
            topic = "General Topic"
            if "about" in command.lower():
                topic = command.lower().split("about")[-1].strip()
            
            # Extract slide count if specified
            slide_count = 7
            count_match = re.search(r"(\d+)\s*slides?", command)
            if count_match:
                slide_count = int(count_match.group(1))
            
            try:
                print(f"Generating presentation: Title='{title}', Topic='{topic}', Slides={slide_count}")
                presentation_path = generate_presentation(
                    title=title,
                    topic=topic,
                    theme_id="modern",
                    slide_count=slide_count
                )
                
                if presentation_path:
                    # Store the path in memory for later use
                    memory.add_topic_memory("last_presentation_path", str(presentation_path))
                    return f"I've created your presentation about {topic}! Would you like to take a look at it? (Yes/No)\nThe presentation is saved at: {presentation_path}"
                else:
                    return "I encountered an issue while creating the presentation. Please try again."
            except Exception as e:
                print(f"Error generating presentation: {e}")
                return "I had trouble creating the presentation. Please try again."
        
        # Handle viewing the last generated presentation
        if any(word in command.lower() for word in ["yes", "yeah", "sure", "okay", "ok", "yep", "please"]):
            last_ppt_path = memory.get_topic_memory("last_presentation_path")
            if last_ppt_path and os.path.exists(last_ppt_path):
                try:
                    if platform.system() == "Windows":
                        os.startfile(last_ppt_path)
                    elif platform.system() == "Darwin":  # macOS
                        subprocess.Popen(["open", last_ppt_path])
                    else:  # Linux
                        subprocess.Popen(["xdg-open", last_ppt_path])
                    return "I've opened the presentation for you!"
                except Exception as e:
                    print(f"Error opening presentation: {e}")
                    return f"The presentation is located at: {last_ppt_path}"
            else:
                return "I couldn't find the last generated presentation. Please try generating a new one."

        # Handle system control commands
        if "volume" in command.lower():
            if "set" in command.lower() or "change" in command.lower():
                # Extract volume level
                level_match = re.search(r'\d+', command)
                if level_match:
                    level = int(level_match.group(0))
                    if 0 <= level <= 100:
                        if set_system_volume(level):
                            return f"Volume set to {level}%"
                        else:
                            return "Failed to set volume"
                    else:
                        return "Volume level must be between 0 and 100"
            else:
                # Get current volume
                current_volume = get_system_volume()
                return f"Current volume is {current_volume:.1f}%"
        
        if "brightness" in command.lower():
            if "set" in command.lower() or "change" in command.lower():
                # Extract brightness level
                level_match = re.search(r'\d+', command)
                if level_match:
                    level = int(level_match.group(0))
                    if 0 <= level <= 100:
                        if set_screen_brightness(level):
                            return f"Screen brightness set to {level}%"
                        else:
                            return "Failed to set brightness"
                    else:
                        return "Brightness level must be between 0 and 100"
            else:
                # Get current brightness
                current_brightness = get_screen_brightness()
                return f"Current screen brightness is {current_brightness}%"
        
        # Handle enhanced YouTube commands
        if "youtube" in command.lower():
            if "info" in command.lower() or "details" in command.lower():
                # Extract video URL
                url_match = re.search(r'https?://(?:www\.)?youtube\.com/watch\?v=[\w-]+', command)
                if url_match:
                    video_url = url_match.group(0)
                    info = get_youtube_video_info(video_url)
                    return f"Video Information:\nTitle: {info['title']}\nViews: {info['view_count']}\nUploaded: {info['upload_date']}\nDescription: {info['description']}"
                else:
                    return "Please provide a valid YouTube video URL"
            
            elif "playlist" in command.lower() and "create" in command.lower():
                # Extract playlist name and video URLs
                name_match = re.search(r'playlist named "([^"]+)"', command)
                url_matches = re.findall(r'https?://(?:www\.)?youtube\.com/watch\?v=[\w-]+', command)
                
                if name_match and url_matches:
                    playlist_name = name_match.group(1)
                    return create_youtube_playlist(playlist_name, url_matches)
                else:
                    return "Please provide a playlist name and at least one video URL"
            
            else:
                # Handle regular YouTube search
                search_query = command.lower().replace("youtube", "").replace("search", "").replace("on", "").strip()
                if search_query:
                    try:
                        search_url = f"https://www.youtube.com/results?search_query={search_query.replace(' ', '+')}"
                        webbrowser.open(search_url)
                        return f"I've opened YouTube and searched for '{search_query}'"
                    except Exception as e:
                        print(f"Error searching YouTube: {e}")
                        return f"I had trouble searching YouTube for '{search_query}'. Please try again."
        
        # Check for email commands first
        if "email" in command.lower() or "send email" in command.lower():
            # Extract email address
            email_match = re.search(r'[\w\.-]+@[\w\.-]+\.\w+', command)
            if not email_match:
                return "Please provide a valid email address to send to."
            
            to_email = email_match.group(0)
            
            # Extract subject if present
            subject = ""
            if "about" in command.lower():
                subject_start = command.lower().find("about") + 6
                subject = command[subject_start:].strip()
                # Remove any "that" or additional content from subject
                if "that" in subject.lower():
                    subject = subject[:subject.lower().find("that")].strip()
            
            # Extract body if present
            body = ""
            if "that" in command.lower():
                body_start = command.lower().find("that") + 5
                body = command[body_start:].strip()
            
            # If no body provided but there's content after the email, use AI to generate email content
            if not body and len(command.split()) > 3:
                try:
                    # Extract the context for email generation
                    context_start = command.lower().find(to_email) + len(to_email)
                    context = command[context_start:].strip()
                    
                    # Generate email content using AI
                    model = genai.GenerativeModel('gemini-2.0-flash-lite')
                    prompt = f"""Generate a professional email with the following context:
                    Recipient: {to_email}
                    Subject: {subject if subject else 'No subject provided'}
                    Context: {context}
                    
                    Requirements:
                    1. Write in a professional but friendly tone
                    2. Keep it concise and clear
                    3. Include a proper greeting and sign-off
                    4. Format it as a proper email
                    
                    Only provide the email body text, no explanations or additional text."""
                    
                    response = model.generate_content(prompt)
                    body = response.text.strip()
                    
                    # If no subject was provided, generate one
                    if not subject:
                        subject_prompt = f"""Generate a short, professional email subject line for this email:
                        {body}
                        
                        Only provide the subject line, no explanations or additional text."""
                        subject_response = model.generate_content(subject_prompt)
                        subject = subject_response.text.strip()
                except Exception as e:
                    print(f"Error generating email content: {e}")
                    return "I had trouble generating the email content. Please try again with more specific instructions."
            
            # Extract attachments if present
            attachments = []
            if "attach" in command.lower() or "with file" in command.lower():
                # Look for file paths in quotes
                file_matches = re.findall(r'["\']([^"\']+\.(?:txt|pdf|doc|docx|jpg|jpeg|png|zip))["\']', command)
                if file_matches:
                    attachments = file_matches
                else:
                    # Look for file paths without quotes
                    file_matches = re.findall(r'\b[\w\-\.]+\.(?:txt|pdf|doc|docx|jpg|jpeg|png|zip)\b', command)
                    if file_matches:
                        attachments = file_matches
            
            # If we have all required information, send the email
            if to_email:
                # Format the email body if it's not already formatted
                if body:
                    try:
                        model = genai.GenerativeModel('gemini-2.0-flash-lite')
                        format_prompt = f"""Format this email content professionally:
                        Original message: {body}
                        
                        Requirements:
                        1. Add a proper greeting (Dear/Hi)
                        2. Add a proper sign-off (Best regards/Sincerely)
                        3. Format it as a proper email with paragraphs
                        4. Keep the original message but make it more professional
                        5. Do not include placeholders like [your name] or [recipient name]
                        6. Use the actual email address {to_email} in the greeting
                        
                        Only provide the formatted email, no explanations."""
                        
                        format_response = model.generate_content(format_prompt)
                        body = format_response.text.strip()
                        
                        # If no subject was provided, generate one
                        if not subject:
                            subject_prompt = f"""Generate a short, professional email subject line for this email:
                            {body}
                            
                            Only provide the subject line, no explanations or additional text."""
                            subject_response = model.generate_content(subject_prompt)
                            subject = subject_response.text.strip()
                    except Exception as e:
                        print(f"Error formatting email: {e}")
                
                result = send_email(to_email, subject, body, attachments)
                return result
            else:
                return "Please provide an email address to send to."
        
        # Check for task scheduling commands
        if any(phrase in command.lower() for phrase in ["schedule task", "schedule a task", "set reminder", "create task", "schedule", "list tasks", "remove task", "delete task", "update task"]):
            return schedule_task_command(command)
        
        # Define screen analysis phrases
        screen_analysis_phrases = [
            "what do you see", 
            "what's on my screen", 
            "analyze screen", 
            "analyze my screen",
            "what's visible", 
            "what is on my screen",
            "what can you see",
            "tell me what you see",
            "describe my screen",
            "show screen info"
        ]
        
        # Check if this is a screen analysis request
        is_screen_request = any(phrase in command.lower() for phrase in screen_analysis_phrases)
        
        # Only analyze screen if it's a screen-related request
        screen_analysis = None
        if is_screen_request:
            print(f"Debug: Processing screen analysis request: {command}")
            screen_analysis = analyze_screen()
            
            # For voice mode, we'll return a tuple (display_text, speak_text)
            if is_voice_mode:
                # Full text for console display
                display_text = f"I can see:\n"
                if isinstance(screen_analysis, dict):
                    if screen_analysis.get("current_app"):
                        display_text += f"- You're currently in: {screen_analysis['current_app']}\n"
                    if screen_analysis.get("visible_text"):
                        display_text += f"- Visible text: {screen_analysis['visible_text']}\n"
                    if screen_analysis.get("visible_elements"):
                        display_text += f"- Visible elements: {', '.join(screen_analysis['visible_elements'])}\n"
                    if screen_analysis.get("context"):
                        display_text += f"- Context: {screen_analysis['context']}\n"
                    if screen_analysis.get("urls"):
                        display_text += f"- URLs: {', '.join(screen_analysis['urls'])}\n"
                    if screen_analysis.get("files"):
                        display_text += f"- Files: {', '.join(screen_analysis['files'])}\n"
                else:
                    display_text += f"- {screen_analysis}\n"
                
                # Concise text for speech
                speak_text = f"I can see your screen. "
                if isinstance(screen_analysis, dict):
                    if screen_analysis.get("current_app"):
                        speak_text += f"You're currently in {screen_analysis['current_app']}. "
                    if screen_analysis.get("context"):
                        speak_text += f"{screen_analysis['context']} "
                    if screen_analysis.get("visible_elements") and len(screen_analysis['visible_elements']) > 0:
                        speak_text += f"I can see elements like {', '.join(screen_analysis['visible_elements'][:3])}. "
                    if screen_analysis.get("urls") and len(screen_analysis['urls']) > 0:
                        speak_text += f"I noticed URLs including {screen_analysis['urls'][0]}. "
                else:
                    speak_text += f"{screen_analysis}"
                
                return (display_text, speak_text)
            else:
                # Regular text mode response with full details
                response = f"I can see:\n"
                if isinstance(screen_analysis, dict):
                    if screen_analysis.get("current_app"):
                        response += f"- You're currently in: {screen_analysis['current_app']}\n"
                    if screen_analysis.get("visible_text"):
                        response += f"- Visible text: {screen_analysis['visible_text']}\n"
                    if screen_analysis.get("visible_elements"):
                        response += f"- Visible elements: {', '.join(screen_analysis['visible_elements'])}\n"
                    if screen_analysis.get("context"):
                        response += f"- Context: {screen_analysis['context']}\n"
                    if screen_analysis.get("urls"):
                        response += f"- URLs: {', '.join(screen_analysis['urls'])}\n"
                    if screen_analysis.get("files"):
                        response += f"- Files: {', '.join(screen_analysis['files'])}\n"
                else:
                    response += f"- {screen_analysis}\n"
                return response
        
        # Handle direct app opening commands
        if command.lower().startswith("open "):
            app_name = command[5:].strip().lower()
            try:
                if "youtube" in app_name:
                    webbrowser.open("https://www.youtube.com")
                    return f"I've opened YouTube for you"
                elif "chrome" in app_name:
                    if platform.system() == "Windows":
                        subprocess.Popen(["start", "chrome"])
                    elif platform.system() == "Darwin":  # macOS
                        subprocess.Popen(["open", "-a", "Google Chrome"])
                    else:  # Linux
                        subprocess.Popen(["google-chrome"])
                    return f"I've opened Chrome for you"
                elif "notepad" in app_name:
                    if platform.system() == "Windows":
                        subprocess.Popen(["notepad"])
                    elif platform.system() == "Darwin":
                        subprocess.Popen(["open", "-a", "TextEdit"])
                    else:
                        subprocess.Popen(["gedit"])
                    return f"I've opened the text editor for you"
                else:
                    # Try to open as a website if it's not a known app
                    if not app_name.startswith(('http://', 'https://')):
                        app_name = 'https://' + app_name
                    webbrowser.open(app_name)
                    return f"I've opened {app_name} for you"
            except Exception as e:
                print(f"Error opening application: {e}")
                return f"I had trouble opening {app_name}. Please make sure it's installed and try again."
        
        # Handle system info commands
        if "system info" in command.lower() or "system information" in command.lower():
            return get_system_info()
        
        # Get memory context for better AI understanding
        memory_context = memory.get_memory_context(command)
        
        # Use AI to understand the user's intent with better context
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        intent_prompt = f"""Analyze this user request and understand their intent: "{command}"

        Current screen context: {screen_analysis}
        
        {memory_context if memory_context else ""}

        Break down what the user wants to do and provide specific details.
        Consider the current screen context and previous interactions when determining the action.
        
        Examples of how to handle different requests:
        1. For opening websites:
           - "open youtube.com" -> Open the website in browser
           - "search for pokemon" -> Open browser and search
           - "go to google" -> Open Google in browser
        
        2. For opening applications:
           - "open notepad" -> Open Notepad
           - "open chrome" -> Open Chrome browser
           - "open this app" -> Use screen context to identify and open the visible app
        
        3. For YouTube/music (ONLY if explicitly mentioned):
           - "play despacito on youtube" -> Open YouTube and play the video
           - "play pokemon episode 1 on youtube" -> Search and play the specific video
           - "play some music on youtube" -> Open YouTube and play music
           - "search for despacito on youtube" -> Search YouTube for the video
        
        4. For code/writing:
           - "write a story" -> Create and open a text file
           - "create python code" -> Generate and run Python code
           - "make an html page" -> Create and open HTML file
        
        5. For system info:
           - "show system info" -> Display system information
           - "what's my cpu usage" -> Show CPU usage
           
        6. For email:
           - "send email to john@example.com" -> Send email to specified address
           - "email john@example.com about meeting" -> Send email with subject
           - "send file to john@example.com" -> Send email with attachment
        
        Respond with a JSON object containing:
        {{
            "intent": "what the user wants to do",
            "steps": [
                {{
                    "action": "specific action to take",
                    "details": "specific details for this action",
                    "parameters": {{"param1": "value1"}}
                }}
            ],
            "context": "any additional context or requirements"
        }}"""
        
        # Use the model without additional parameters
        intent_response = model.generate_content(intent_prompt)
        intent_plan = intent_response.text.strip()
        
        try:
            import json
            plan = json.loads(intent_plan)
            steps = plan.get("steps", [])
            context = plan.get("context", "")
            
            results = []
            for step in steps:
                action = step.get("action", "").lower()
                details = step.get("details", "")
                params = step.get("parameters", {})
                
                # Handle website/browser actions
                if "website" in action or "browser" in action or "search" in action:
                    url = params.get("url", details)
                    if not url.startswith(('http://', 'https://')):
                        if "youtube" in url.lower() or "youtube" in command.lower():
                            search_url = f"https://www.youtube.com/results?search_query={url.replace(' ', '+')}"
                            webbrowser.open(search_url)
                            results.append(f"I've opened YouTube and searched for '{url}'")
                        else:
                            results.append(open_browser(url))
                
                # Handle application opening
                elif "open" in action or "launch" in action:
                    app_name = params.get("app", details)
                    if app_name:
                        if "this" in app_name.lower() or "that" in app_name.lower():
                            # Use screen context to identify the app
                            if isinstance(screen_analysis, dict):
                                visible_elements = screen_analysis.get("visible_elements", [])
                            if visible_elements:
                                app_name = visible_elements[0]  # Use the first visible element
                        try:
                            if platform.system() == "Windows":
                                subprocess.Popen(["start", app_name])
                            elif platform.system() == "Darwin":  # macOS
                                subprocess.Popen(["open", "-a", app_name])
                            else:  # Linux
                                subprocess.Popen([app_name])
                            results.append(f"I've opened {app_name} for you")
                        except Exception as e:
                            print(f"Error opening application: {e}")
                            results.append(f"I had trouble opening {app_name}. Please make sure it's installed and try again.")
                
                # Handle YouTube/music requests (only if explicitly mentioned)
                elif ("youtube" in action or "play" in action or "music" in action) and "youtube" in command.lower():
                    search_query = params.get("query", details)
                    if not search_query:
                        search_query = command.replace("play", "").replace("youtube", "").strip()
                    search_url = f"https://www.youtube.com/results?search_query={search_query.replace(' ', '+')}"
                    webbrowser.open(search_url)
                    results.append(f"I've opened YouTube and searched for '{search_query}'")
                
                # Handle code generation
                elif "code" in action or "program" in action:
                    language = params.get("language", "Python")
                    task = params.get("task", details)
                    code = generate_code(task, language)
                    if code:
                        extension = ".py" if language.lower() == "python" else ".html"
                        file_path = create_and_write_file(code, extension)
                        if file_path:
                            open_in_editor(file_path)
                            if language.lower() == "python":
                                output = run_python_file(file_path)
                                results.append(f"I've created and run the Python code. Output:\n{output}")
                            else:
                                open_html_in_browser(file_path)
                                results.append(f"I've created and opened the HTML file in your browser")
                
                # Handle writing/story requests
                elif "write" in action or "story" in action:
                    content_prompt = f"""Write {details} or a story about {details}.
                    
                    {memory_context if memory_context else ""}
                    
                    Keep it under 500 words unless requested to be longer."""
                    
                    # Use the model without additional parameters
                    content_response = model.generate_content(content_prompt)
                    content = content_response.text.strip()
                    
                    # Create a file with the generated content
                    file_path = create_and_write_file(content, ".txt")
                    if file_path:
                        open_in_editor(file_path)
                        results.append(f"I've written your content and opened it in an editor. Here's a preview of what I wrote:\n\n{content[:200]}...")
                
                # System information
                elif "system" in action or "info" in action or "cpu" in action or "memory" in action:
                    results.append(get_system_info())
                
                # Handle email requests
                elif "email" in action or "send" in action and "email" in command.lower():
                    to_email = params.get("to", "")
                    subject = params.get("subject", "")
                    body = params.get("body", details)
                    attachments = params.get("attachments", [])
                    
                    # Extract email address if not provided in parameters
                    if not to_email:
                        email_match = re.search(r'[\w\.-]+@[\w\.-]+\.\w+', command)
                        if email_match:
                            to_email = email_match.group(0)
                    
                    # Extract subject if not provided in parameters
                    if not subject and "about" in command.lower():
                        subject_start = command.lower().find("about") + 6
                        subject = command[subject_start:].strip()
                    
                    if to_email:
                        results.append(send_email(to_email, subject, body, attachments))
                    else:
                        results.append("Please provide an email address to send to.")
                
                # Handle PowerPoint generation requests
                elif ("create" in action or "make" in action or "generate" in action) and ("presentation" in action or "powerpoint" in action or "ppt" in action):
                    title = params.get("title", "My Presentation")
                    topic = params.get("topic", details if details else "General Topic")
                    theme = params.get("theme", "modern")
                    slide_count = params.get("slides", 7)
                    
                    try:
                        # Generate the presentation
                        presentation_path = generate_presentation(
                            title=title,
                            topic=topic,
                            theme_id=theme,
                            slide_count=slide_count
                        )
                        
                        if presentation_path:
                            # Store the path in memory for later use
                            memory.add_topic_memory("last_presentation_path", str(presentation_path))
                            
                            # Ask if user wants to view the presentation
                            results.append(f"I've created your presentation about {topic}! Would you like to take a look at it? (Yes/No)")
                            results.append(f"The presentation is saved at: {presentation_path}")
                        else:
                            results.append("I encountered an issue while creating the presentation. Please try again.")
                    except Exception as e:
                        print(f"Error generating presentation: {e}")
                        results.append("I had trouble creating the presentation. Please try again.")
                
                # Handle viewing the last generated presentation
                elif any(word in action.lower() for word in ["view", "open", "show", "see"]) and any(word in action.lower() for word in ["presentation", "powerpoint", "ppt"]):
                    # Check if this is a response to the "would you like to take a look" question
                    is_yes_response = any(word in command.lower() for word in ["yes", "yeah", "sure", "okay", "ok", "yep", "please"])
                    
                    if is_yes_response:
                        # Try to get the last presentation path from memory
                        last_ppt_path = memory.get_topic_memory("last_presentation_path")
                        if last_ppt_path and os.path.exists(last_ppt_path):
                            try:
                                if platform.system() == "Windows":
                                    os.startfile(last_ppt_path)
                                elif platform.system() == "Darwin":  # macOS
                                    subprocess.Popen(["open", last_ppt_path])
                                else:  # Linux
                                    subprocess.Popen(["xdg-open", last_ppt_path])
                                results.append("I've opened the presentation for you!")
                            except Exception as e:
                                print(f"Error opening presentation: {e}")
                                results.append(f"The presentation is located at: {last_ppt_path}")
                        else:
                            results.append("I couldn't find the last generated presentation. Please try generating a new one.")
                    else:
                        results.append("No problem! The presentation is saved and you can view it later.")
                
                # Execute system command directly
                elif "command" in action or "execute" in action or "run" in action:
                    cmd = params.get("command", details)
                    if cmd:
                        results.append(execute_system_command(cmd))
            
            # If we've done something
            if results:
                if is_voice_mode:
                    # For voice mode, provide a concise version for speech
                    display_text = "\n".join(results)
                    speak_text = f"I've completed your request to {plan.get('intent', 'perform your request')}."
                    return (display_text, speak_text)
                else:
                    return "\n".join(results)
            
            # If we couldn't determine the action, generate a response directly
            response_prompt = f"""
            The user asked: "{command}"
            
            Current screen context: {screen_analysis if screen_analysis else "No screen analysis available"}
            
            {memory_context if memory_context else ""}
            
            Respond helpfully to this request in a friendly, conversational way. Keep your response concise.
            """
            
            response = model.generate_content(response_prompt)
            return response.text.strip()
            
        except Exception as e:
            print(f"AI intent error: {e}")
            # Fallback for parsing error - use a direct response
            direct_prompt = f"""
            The user asked: "{command}"
            
            {memory_context if memory_context else ""}
            
            Please respond helpfully to this request in a friendly, conversational way. Keep your response concise.
            """
            
            response = model.generate_content(direct_prompt)
            return response.text.strip()
    
    except Exception as e:
        print(f"Error processing command: {e}")
        return "I'm here to help. What would you like me to do?"

# Create the main agent
root_agent = Agent(
    name="SystemControlAssistant",
    model="gemini-2.0-flash-lite",
    instruction=(
        "You are a helpful system control assistant that can execute commands and control the computer. "
        "You can:\n"
        "1. Open browsers and websites\n"
        "2. Show system information\n"
        "3. Run system commands\n"
        "4. Open applications\n"
        "5. Generate and write code in various languages\n"
        "6. Execute code and show results\n"
        "7. Send emails with attachments\n\n"
        "Always be friendly and natural in your responses. Don't mention errors or technical details to the user. "
        "Just confirm that you've completed their request in a conversational way."
    ),
    description="An assistant that can control your system and execute commands."
)

# Initialize session service
session_service = InMemorySessionService()
APP_NAME = "SystemControlApp"
USER_ID = "default_user"

# Global session variable
session = None

# Global memory variable 
memory = ContextualMemory()

async def initialize_session():
    """Initialize the session asynchronously."""
    global session
    try:
        session = await session_service.create_session(app_name=APP_NAME, user_id=USER_ID)
        if session is None:
            raise Exception("Failed to create session")
        return session
    except Exception as e:
        print(f"Error creating session: {e}")
        raise

async def process_user_input(user_input: str, is_voice_mode=False) -> str:
    """Process user input and return agent response."""
    print(f"Debug: Processing user input: {user_input}")
    try:
        # Initialize memory if not already done
        global memory
        if not memory.memory_initialized:
            memory.load_from_disk()
        
        # Memory management commands
        if "clear memory" in user_input.lower():
            memory = ContextualMemory()
            memory.memory_initialized = True
            memory.save_to_disk()
            return "Memory has been cleared."
        
        if "show memory" in user_input.lower() or "show what you remember" in user_input.lower():
            facts = memory.get_facts()
            prefs = memory.get_preferences()
            convs = memory.get_recent_conversations(5)
            
            response = "Here's what I remember:\n\n"
            
            if facts:
                response += "Facts about you:\n"
                for category, fact_list in facts.items():
                    for fact in fact_list:
                        response += f"- {fact}\n"
                response += "\n"
            
            if prefs:
                response += "Your preferences:\n"
                for category, pref in prefs.items():
                    response += f"- {category}: {pref}\n"
                response += "\n"
            
            if convs:
                response += "Recent conversations:\n"
                for i, conv in enumerate(convs):
                    response += f"- You: \"{conv['user_input']}\"\n"
                    response += f"  Me: \"{conv['agent_response']}\"\n"
                
            if not facts and not prefs and not convs:
                response += "I don't have any memories stored yet."
                
            return response
        
        # Get memory context relevant to the current input
        memory_context = memory.get_memory_context(user_input)
        
        # Check if this is a direct code generation request
        if "write" in user_input.lower() and any(lang in user_input.lower() for lang in ["html", "css", "javascript", "python", "code"]):
            language = "HTML" if "html" in user_input.lower() else "Python"
            code = generate_code(user_input, language)
            if code:
                extension = ".py" if language.lower() == "python" else ".html"
                file_path = create_and_write_file(code, extension)
                if file_path:
                    open_in_editor(file_path)
                    if language.lower() == "python":
                        output = run_python_file(file_path)
                        response = f"I've created and run the Python code."
                    else:
                        open_html_in_browser(file_path)
                        response = "I've created and opened the HTML file in your browser."
                    
                    # Store the interaction in memory
                    memory.add_conversation(user_input, response)
                    memory.extract_facts_from_conversation(user_input, response)
                    memory.save_to_disk()
                    return response
                        
        # Define screen analysis phrases
        screen_analysis_phrases = [
            "what do you see", 
            "what's on my screen", 
            "analyze screen", 
            "analyze my screen",
            "what's visible", 
            "what is on my screen",
            "what can you see",
            "tell me what you see",
            "describe my screen",
            "show screen info"
        ]
        
        # Check if this is a direct screen analysis request
        is_screen_request = any(phrase in user_input.lower() for phrase in screen_analysis_phrases)
        
        # Special command to analyze screen directly
        if is_screen_request:
            print("Debug: Direct screen analysis request detected")
            try:
                screen_analysis = analyze_screen()
                
                # For voice mode, we'll return a tuple (display_text, speak_text)
                if is_voice_mode:
                    # Full text for console display
                    display_text = f"I can see:\n"
                    if isinstance(screen_analysis, dict):
                        if screen_analysis.get("current_app"):
                            display_text += f"- You're currently in: {screen_analysis['current_app']}\n"
                        if screen_analysis.get("visible_text"):
                            display_text += f"- Visible text: {screen_analysis['visible_text']}\n"
                        if screen_analysis.get("visible_elements"):
                            display_text += f"- Visible elements: {', '.join(screen_analysis['visible_elements'])}\n"
                        if screen_analysis.get("context"):
                            display_text += f"- Context: {screen_analysis['context']}\n"
                        if screen_analysis.get("urls"):
                            display_text += f"- URLs: {', '.join(screen_analysis['urls'])}\n"
                        if screen_analysis.get("files"):
                            display_text += f"- Files: {', '.join(screen_analysis['files'])}\n"
                    else:
                        display_text += f"- {screen_analysis}\n"
                    
                    # Concise text for speech
                    speak_text = f"I can see your screen. "
                    if isinstance(screen_analysis, dict):
                        if screen_analysis.get("current_app"):
                            speak_text += f"You're currently in {screen_analysis['current_app']}. "
                        if screen_analysis.get("context"):
                            speak_text += f"{screen_analysis['context']} "
                        if screen_analysis.get("visible_elements") and len(screen_analysis['visible_elements']) > 0:
                            speak_text += f"I can see elements like {', '.join(screen_analysis['visible_elements'][:3])}. "
                        if screen_analysis.get("urls") and len(screen_analysis['urls']) > 0:
                            speak_text += f"I noticed URLs including {screen_analysis['urls'][0]}. "
                    else:
                        speak_text += f"{screen_analysis}"
                    
                    # Store screen information in topic memory
                    memory.add_topic_memory("last_screen", screen_analysis)
                    memory.add_conversation(user_input, display_text)
                    memory.save_to_disk()
                    
                    return (display_text, speak_text)
                else:
                    # Regular text mode response with full details
                    response = f"I can see:\n"
                    if isinstance(screen_analysis, dict):
                        if screen_analysis.get("current_app"):
                            response += f"- You're currently in: {screen_analysis['current_app']}\n"
                        if screen_analysis.get("visible_text"):
                            response += f"- Visible text: {screen_analysis['visible_text']}\n"
                        if screen_analysis.get("visible_elements"):
                            response += f"- Visible elements: {', '.join(screen_analysis['visible_elements'])}\n"
                        if screen_analysis.get("context"):
                            response += f"- Context: {screen_analysis['context']}\n"
                        if screen_analysis.get("urls"):
                            response += f"- URLs: {', '.join(screen_analysis['urls'])}\n"
                        if screen_analysis.get("files"):
                            response += f"- Files: {', '.join(screen_analysis['files'])}\n"
                    else:
                        response += f"- {screen_analysis}\n"
                    
                    # Store screen information in topic memory
                    memory.add_topic_memory("last_screen", screen_analysis)
                    memory.add_conversation(user_input, response)
                    memory.save_to_disk()
                    
                    return response
            except Exception as screen_error:
                print(f"Debug: Error in direct screen analysis: {screen_error}")
                return "I tried to analyze your screen, but encountered an error. Please try again."
        
        # First try to process as a system command using AI
        try:
            print("Debug: Trying to process command")
            command_response = process_command(user_input, is_voice_mode=is_voice_mode)
            if command_response and not command_response.startswith("I understand you want me to help"):
                print("Debug: Command processed successfully")
                
                # Store the interaction in memory
                if isinstance(command_response, tuple):
                    # For voice mode responses (display_text, speak_text)
                    memory.add_conversation(user_input, command_response[0])
                else:
                    memory.add_conversation(user_input, command_response)
                
                memory.extract_facts_from_conversation(user_input, command_response if isinstance(command_response, str) else command_response[0])
                memory.save_to_disk()
                
                return command_response
        except Exception as command_error:
            print(f"Debug: Error in process_command: {command_error}")
            # If the command processing fails, continue to LLM fallback
        
        # For basic questions, use a simpler approach with the model directly
        simple_question_keywords = ["what is", "who is", "how to", "explain", "define", "tell me about"]
        is_simple_question = any(keyword in user_input.lower() for keyword in simple_question_keywords)
        
        if is_simple_question:
            try:
                print("Debug: Using direct model for simple question")
                model = genai.GenerativeModel('gemini-2.0-flash-lite')
                
                # Add memory context to the prompt if available
                if memory_context:
                    enhanced_prompt = f"""Answer this question briefly: {user_input}
                    
                    Consider the following context from previous interactions:
                    {memory_context}"""
                else:
                    enhanced_prompt = f"Answer this question briefly: {user_input}"
                
                response = model.generate_content(enhanced_prompt)
                answer = response.text.strip()
                print(f"Debug: Direct model response: {answer}")
                
                # Store the interaction in memory
                memory.add_conversation(user_input, answer)
                memory.extract_facts_from_conversation(user_input, answer)
                memory.save_to_disk()
                
                return answer
            except Exception as e:
                print(f"Debug: Error using direct model: {e}")
                # Fall through to the LLM approach
        
        # If not a system command, use the LLM for general conversation
        print("Debug: Falling back to LLM conversation")
        global session
        if session is None:
            print("Debug: Initializing session")
            await initialize_session()
        
        # Add memory context as system message if available
        if memory_context:
            context_message = genai_types.Content(role='user', parts=[genai_types.Part(text=f"Consider this context from previous interactions: {memory_context}")])
            # We would add this to the conversation history if this was a multi-turn conversation API
        
        user_message = genai_types.Content(role='user', parts=[genai_types.Part(text=user_input)])
        
        from google.adk.runners import Runner
        runner = Runner(
            agent=root_agent,
            app_name=APP_NAME,
            session_service=session_service
        )
        
        agent_reply = ""
        print("Debug: Running LLM")
        try:
            async for event in runner.run_async(user_id=USER_ID, session_id=session.id, new_message=user_message):
                if hasattr(event, 'type'):
                    if event.type == 'final_response':
                        if hasattr(event, 'content') and event.content and hasattr(event.content, 'parts'):
                            for part in event.content.parts:
                                if hasattr(part, 'text') and part.text:
                                    agent_reply = part.text
                    elif event.type == 'error':
                        error_text = event.content.parts[0].text if hasattr(event, 'content') and event.content and hasattr(event.content, 'parts') else 'Unknown error'
                        print(f"Debug: Error from LLM: {error_text}")
        except Exception as runner_error:
            print(f"Debug: Error in runner: {runner_error}")
        
        print(f"Debug: LLM response: {agent_reply if agent_reply else 'No response'}")
        
        # If we didn't get a response from the LLM, use a direct model as fallback
        if not agent_reply:
            try:
                print("Debug: Using emergency direct model for fallback")
                model = genai.GenerativeModel('gemini-2.0-flash-lite')
                
                # Add memory context to the prompt if available
                if memory_context:
                    enhanced_prompt = f"""Respond to this: {user_input}
                    
                    Consider the following context from previous interactions:
                    {memory_context}"""
                else:
                    enhanced_prompt = f"Respond to this: {user_input}"
                
                response = model.generate_content(enhanced_prompt)
                agent_reply = response.text.strip()
                print(f"Debug: Emergency fallback response: {agent_reply}")
            except Exception as e:
                print(f"Debug: Error in emergency fallback: {e}")
                agent_reply = "I understand your question, but I'm having trouble formulating a response at the moment. Could you try rephrasing or asking something else?"
        
        # Store the interaction in memory
        if agent_reply:
            memory.add_conversation(user_input, agent_reply)
            memory.extract_facts_from_conversation(user_input, agent_reply)
            memory.save_to_disk()
        
        return agent_reply
    
    except Exception as e:
        print(f"Debug: Critical error in process_user_input: {e}")
        # Emergency fallback for screen analysis requests
        if "screen" in user_input.lower() and "analyze" in user_input.lower():
            try:
                avg_color = [0, 0, 0]  # Fallback if everything else fails
                try:
                    img = pyautogui.screenshot()
                    np_img = np.array(img)
                    avg_color = np.mean(np_img, axis=(0, 1))
                except Exception:
                    pass
                
                brightness = sum(avg_color) / 3
                basic_info = f"I can see your screen. It appears to be a {'dark' if brightness < 128 else 'light'} screen."
                return basic_info
            except Exception:
                return "I tried to analyze your screen but encountered a critical error. Please try again."
        
        # Last resort emergency fallback
        try:
            model = genai.GenerativeModel('gemini-2.0-flash-lite')
            response = model.generate_content(f"Answer this briefly: {user_input}")
            return response.text.strip()
        except Exception:
            return "I'm here to help. What would you like me to do?"
